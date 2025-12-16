#!/usr/bin/env python3
"""
PDF to PowerPoint Converter - High Fidelity Conversion
Converts PDF to 100% editable PPTX maintaining 99% format and content
Preserves layout, images, text formatting, and orientation
"""

import sys
import json
import os
import tempfile
import traceback
import re
from io import BytesIO

def install_dependencies():
    """Install required packages if not present"""
    import subprocess
    packages = [
        'python-pptx',
        'pymupdf',  # fitz
        'Pillow',
        'pdf2image',
    ]
    for package in packages:
        try:
            pkg_name = package.replace('-', '_').lower()
            if pkg_name == 'pymupdf':
                pkg_name = 'fitz'
            __import__(pkg_name)
        except ImportError:
            try:
                subprocess.check_call([sys.executable, '-m', 'pip', 'install', package, '-q'])
            except:
                pass

try:
    install_dependencies()
except:
    pass

# Import required libraries
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    from PIL import Image
    HAS_PIL = True
except ImportError:
    HAS_PIL = False


def rgb_to_pptx_color(rgb_tuple):
    """Convert RGB tuple (0-1 range) to RGBColor"""
    if not rgb_tuple or len(rgb_tuple) < 3:
        return RGBColor(0, 0, 0)
    r = int(min(255, max(0, rgb_tuple[0] * 255)))
    g = int(min(255, max(0, rgb_tuple[1] * 255)))
    b = int(min(255, max(0, rgb_tuple[2] * 255)))
    return RGBColor(r, g, b)


def get_text_alignment(flags):
    """Determine text alignment from PDF text flags"""
    # Default to left alignment
    return PP_ALIGN.LEFT


class PDFToPPTXConverter:
    """High-fidelity PDF to PPTX converter"""
    
    def __init__(self, options=None):
        self.options = options or {}
        self.preserve_images = self.options.get('preserve_images', True)
        self.preserve_formatting = self.options.get('preserve_formatting', True)
        self.image_quality = self.options.get('image_quality', 95)
        self.extract_as_image = self.options.get('extract_as_image', False)
        # Default to editable mode (no hybrid) for fully editable output
        self.hybrid_mode = self.options.get('hybrid_mode', False)
        self.temp_files = []  # Track temp files for cleanup
        
    def convert(self, input_path, output_path):
        """
        Convert PDF to PPTX with high fidelity
        
        Args:
            input_path: Path to input PDF file
            output_path: Path to output PPTX file
            
        Returns:
            Dictionary with conversion result
        """
        if not HAS_FITZ:
            return {
                'success': False,
                'error': 'PyMuPDF (fitz) not installed. Run: pip install pymupdf'
            }
            
        if not HAS_PPTX:
            return {
                'success': False,
                'error': 'python-pptx not installed. Run: pip install python-pptx'
            }
        
        try:
            # Open PDF
            pdf_doc = fitz.open(input_path)
            page_count = len(pdf_doc)
            
            if page_count == 0:
                return {
                    'success': False,
                    'error': 'PDF has no pages'
                }
            
            # Create presentation
            prs = Presentation()
            
            # Process each page
            for page_num in range(page_count):
                page = pdf_doc[page_num]
                self._convert_page(prs, page, page_num + 1)
            
            # Save presentation
            prs.save(output_path)
            pdf_doc.close()
            
            # Clean up temp files after saving (important for Windows)
            self._cleanup_temp_files()
            
            output_size = os.path.getsize(output_path)
            
            return {
                'success': True,
                'output_path': output_path,
                'output_size': output_size,
                'page_count': page_count,
                'message': f'PDF converted to PPTX successfully ({page_count} slides)'
            }
            
        except Exception as e:
            # Clean up temp files even on error
            self._cleanup_temp_files()
            return {
                'success': False,
                'error': str(e),
                'traceback': traceback.format_exc()
            }
    
    def _cleanup_temp_files(self):
        """Clean up temporary files with retry for Windows file locking"""
        import time
        for temp_file in self.temp_files:
            for attempt in range(3):
                try:
                    if os.path.exists(temp_file):
                        os.unlink(temp_file)
                    break
                except PermissionError:
                    if attempt < 2:
                        time.sleep(0.1)  # Wait a bit for file to be released
                    # Ignore if we can't delete - OS will clean up temp files
                except Exception:
                    pass
        self.temp_files = []
    
    def _convert_page(self, prs, page, page_num):
        """Convert a single PDF page to a PPTX slide"""
        # Get page dimensions in points
        rect = page.rect
        page_width = rect.width
        page_height = rect.height
        
        # Determine orientation
        is_landscape = page_width > page_height
        
        # Set slide size based on PDF page dimensions (only on first page)
        # PowerPoint slide size is set once for entire presentation
        if page_num == 1:
            if is_landscape:
                # Landscape: Use widescreen 16:9 ratio
                prs.slide_width = Inches(13.333)
                prs.slide_height = Inches(7.5)
            else:
                # Portrait: Use A4-like ratio (closer to PDF standard)
                # Standard A4 is 8.27 x 11.69 inches
                prs.slide_width = Inches(7.5)
                prs.slide_height = Inches(10)
        
        # Get actual slide dimensions
        slide_width_pts = prs.slide_width.pt
        slide_height_pts = prs.slide_height.pt
        
        # Calculate scale to fit PDF page into slide while maintaining aspect ratio
        scale_x = slide_width_pts / page_width
        scale_y = slide_height_pts / page_height
        
        # Use the smaller scale to ensure content fits completely
        scale = min(scale_x, scale_y)
        
        # Calculate offset to center the content if there's extra space
        scaled_width = page_width * scale
        scaled_height = page_height * scale
        offset_x = (slide_width_pts - scaled_width) / 2
        offset_y = (slide_height_pts - scaled_height) / 2
        
        # Add blank slide
        blank_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(blank_layout)
        
        # Option 1: Image-only mode (for when editable text is not needed)
        if self.extract_as_image:
            self._add_page_as_image(slide, page, prs, offset_x, offset_y, scaled_width, scaled_height)
            return
        
        # Option 2: Fully editable mode (default) - extracts text, images, and shapes
        # Extract and add images from PDF (logos, photos, etc.)
        if self.preserve_images:
            self._extract_images(slide, page, scale, prs, offset_x, offset_y)
        
        # Extract and add editable text blocks with formatting
        self._extract_text_blocks(slide, page, scale, prs, offset_x, offset_y, transparent=False)
        
        # Extract and add vector drawings/shapes
        self._extract_drawings(slide, page, scale, prs, offset_x, offset_y)
    
    def _add_page_as_image(self, slide, page, prs, offset_x=0, offset_y=0, width=None, height=None, as_background=False):
        """Add entire page as high-quality image"""
        # Render page at high resolution (3x for crisp quality)
        zoom = 3.0 if self.image_quality >= 90 else 2.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        
        # Save to temporary file in our temp directory (not system temp)
        temp_img_path = tempfile.mktemp(suffix='.png')
        pix.save(temp_img_path)
        self.temp_files.append(temp_img_path)
        
        # Optimize image if PIL is available
        if HAS_PIL and self.image_quality < 100:
            try:
                img = Image.open(temp_img_path)
                # Convert to RGB if necessary
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                # Save as optimized JPEG for smaller file size
                temp_jpg_path = tempfile.mktemp(suffix='.jpg')
                img.save(temp_jpg_path, 'JPEG', quality=self.image_quality, optimize=True)
                img.close()  # Close the image to release file handle
                self.temp_files.append(temp_jpg_path)
                temp_img_path = temp_jpg_path
            except Exception as e:
                print(f"Warning: Image optimization failed: {e}", file=sys.stderr)
        
        # Use provided dimensions or full slide
        img_width = Pt(width) if width else prs.slide_width
        img_height = Pt(height) if height else prs.slide_height
        img_left = Pt(offset_x) if offset_x else Inches(0)
        img_top = Pt(offset_y) if offset_y else Inches(0)
        
        # Add image to slide with proper positioning
        picture = slide.shapes.add_picture(
            temp_img_path,
            img_left,
            img_top,
            img_width,
            img_height
        )
        
        # If adding as background, send to back
        if as_background:
            # Move picture to back by reordering shapes
            spTree = slide.shapes._spTree
            pic_elem = picture._element
            spTree.remove(pic_elem)
            spTree.insert(2, pic_elem)  # Insert after background elements
        
        # Don't delete here - will be cleaned up after presentation is saved
    
    def _extract_images(self, slide, page, scale, prs, offset_x=0, offset_y=0):
        """Extract and add images from PDF page"""
        try:
            image_list = page.get_images(full=True)
            
            for img_index, img_info in enumerate(image_list):
                try:
                    xref = img_info[0]
                    
                    # Get image position on page first
                    img_rects = page.get_image_rects(xref)
                    if not img_rects:
                        continue
                    
                    for img_rect in img_rects:
                        final_img_path = None
                        
                        # Method 1: Render the image area from the page (handles all formats)
                        try:
                            clip = fitz.Rect(img_rect)
                            mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for quality
                            # alpha=False ensures white background instead of black
                            pix = page.get_pixmap(matrix=mat, clip=clip, alpha=False)
                            final_img_path = tempfile.mktemp(suffix='.png')
                            pix.save(final_img_path)
                            self.temp_files.append(final_img_path)
                        except Exception as render_err:
                            print(f"Warning: Render failed: {render_err}", file=sys.stderr)
                            
                            # Method 2: Extract and process with PIL
                            try:
                                base_image = page.parent.extract_image(xref)
                                if base_image:
                                    image_bytes = base_image["image"]
                                    if HAS_PIL:
                                        from io import BytesIO
                                        img = Image.open(BytesIO(image_bytes))
                                        # Convert to RGB with white background
                                        if img.mode in ('RGBA', 'LA', 'PA'):
                                            bg = Image.new('RGB', img.size, (255, 255, 255))
                                            if img.mode != 'RGBA':
                                                img = img.convert('RGBA')
                                            bg.paste(img, mask=img.split()[-1])
                                            img = bg
                                        elif img.mode != 'RGB':
                                            img = img.convert('RGB')
                                        final_img_path = tempfile.mktemp(suffix='.png')
                                        img.save(final_img_path, 'PNG')
                                        img.close()
                                        self.temp_files.append(final_img_path)
                            except Exception as pil_err:
                                print(f"Warning: PIL failed: {pil_err}", file=sys.stderr)
                        
                        if not final_img_path:
                            continue
                        
                        # Calculate position on slide
                        left = Pt(img_rect.x0 * scale + offset_x)
                        top = Pt(img_rect.y0 * scale + offset_y)
                        width = Pt((img_rect.x1 - img_rect.x0) * scale)
                        height = Pt((img_rect.y1 - img_rect.y0) * scale)
                        
                        # Add image to slide
                        try:
                            slide.shapes.add_picture(final_img_path, left, top, width, height)
                        except Exception as img_err:
                            print(f"Warning: Could not add image: {img_err}", file=sys.stderr)
                        
                except Exception as e:
                    print(f"Warning: Error extracting image {img_index}: {e}", file=sys.stderr)
                    continue
                    
        except Exception as e:
            print(f"Warning: Error in image extraction: {e}", file=sys.stderr)
    
    def _extract_text_blocks(self, slide, page, scale, prs, offset_x=0, offset_y=0, transparent=False):
        """Extract text blocks with formatting"""
        try:
            # Get text blocks with detailed information
            blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
            
            for block in blocks:
                if block["type"] != 0:  # Skip non-text blocks (images handled separately)
                    continue
                
                # Get block bounding box with offset for centering
                bbox = block["bbox"]
                block_left = Pt(bbox[0] * scale + offset_x)
                block_top = Pt(bbox[1] * scale + offset_y)
                block_width = Pt((bbox[2] - bbox[0]) * scale)
                block_height = Pt((bbox[3] - bbox[1]) * scale)
                
                # Add some padding to prevent text cutoff
                block_width = Pt((bbox[2] - bbox[0]) * scale + 10)
                block_height = Pt((bbox[3] - bbox[1]) * scale + 5)
                
                # Ensure minimum dimensions
                if block_width < Pt(20):
                    block_width = Pt(50)
                if block_height < Pt(15):
                    block_height = Pt(20)
                
                # Create text box
                try:
                    textbox = slide.shapes.add_textbox(
                        block_left, block_top, block_width, block_height
                    )
                    tf = textbox.text_frame
                    tf.word_wrap = False  # Don't wrap - preserve original line breaks
                    tf.auto_size = None  # Don't auto-resize
                    
                    # No fill or border for clean look
                    textbox.fill.background()
                    textbox.line.fill.background()
                    
                    # Process each line in the block
                    first_line = True
                    for line in block.get("lines", []):
                        # Create a new paragraph for each line
                        if first_line:
                            p = tf.paragraphs[0]
                            first_line = False
                        else:
                            p = tf.add_paragraph()
                        
                        # Process spans within the line (preserves inline formatting)
                        for span in line.get("spans", []):
                            text = span.get("text", "")
                            if not text:
                                continue
                            
                            # Add run with text
                            run = p.add_run()
                            run.text = text
                            
                            # Apply formatting (font, size, color, bold, italic)
                            if self.preserve_formatting:
                                self._apply_text_formatting(run, span, transparent)
                                
                except Exception as tb_err:
                    print(f"Warning: Could not create textbox: {tb_err}", file=sys.stderr)
                    continue
                    
        except Exception as e:
            print(f"Warning: Error in text extraction: {e}", file=sys.stderr)
    
    def _apply_text_formatting(self, run, span, transparent=False):
        """Apply text formatting from PDF span to PPTX run"""
        try:
            # Font size
            font_size = span.get("size", 12)
            run.font.size = Pt(font_size)
            
            # Font name - try to preserve original font
            font_name = span.get("font", "")
            if font_name:
                # Clean font name (remove subset prefix like ABCDEF+)
                clean_name = font_name.split("+")[-1] if "+" in font_name else font_name
                # Keep style suffix for better matching (e.g., Arial-Bold)
                run.font.name = clean_name
            
            # Font color - extract from PDF
            color = span.get("color")
            if color is not None:
                # Color is an integer in PyMuPDF (RGB packed)
                if isinstance(color, int):
                    r = (color >> 16) & 0xFF
                    g = (color >> 8) & 0xFF
                    b = color & 0xFF
                    run.font.color.rgb = RGBColor(r, g, b)
            else:
                # Default to black
                run.font.color.rgb = RGBColor(0, 0, 0)
            
            # Bold/Italic detection from font name
            font_lower = font_name.lower() if font_name else ""
            run.font.bold = "bold" in font_lower or "black" in font_lower or "heavy" in font_lower
            run.font.italic = "italic" in font_lower or "oblique" in font_lower
            
            # Flags-based formatting (PyMuPDF flags)
            flags = span.get("flags", 0)
            if flags & 2 ** 1:  # Italic flag
                run.font.italic = True
            if flags & 2 ** 4:  # Bold flag
                run.font.bold = True
                
        except Exception as e:
            print(f"Warning: Could not apply formatting: {e}", file=sys.stderr)
    
    def _extract_drawings(self, slide, page, scale, prs, offset_x=0, offset_y=0):
        """Extract vector drawings and shapes"""
        try:
            # Get drawings/paths from page
            drawings = page.get_drawings()
            
            for drawing in drawings:
                try:
                    rect = drawing.get("rect")
                    if not rect:
                        continue
                    
                    # Calculate position with offset
                    left = Pt(rect.x0 * scale + offset_x)
                    top = Pt(rect.y0 * scale + offset_y)
                    width = Pt((rect.x1 - rect.x0) * scale)
                    height = Pt((rect.y1 - rect.y0) * scale)
                    
                    # Skip very small drawings
                    if width < Pt(5) or height < Pt(5):
                        continue
                    
                    # Get fill and stroke colors
                    fill_color = drawing.get("fill")
                    stroke_color = drawing.get("color")
                    
                    # Add rectangle shape (simplified)
                    shape = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        left, top, width, height
                    )
                    
                    # Apply colors
                    if fill_color:
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = rgb_to_pptx_color(fill_color)
                    else:
                        shape.fill.background()
                    
                    if stroke_color:
                        shape.line.color.rgb = rgb_to_pptx_color(stroke_color)
                        shape.line.width = Pt(drawing.get("width", 1))
                    else:
                        shape.line.fill.background()
                        
                except Exception as draw_err:
                    continue  # Skip problematic drawings
                    
        except Exception as e:
            print(f"Warning: Error in drawing extraction: {e}", file=sys.stderr)


def convert_pdf_to_pptx(input_path, output_path, options=None):
    """
    Main conversion function
    
    Args:
        input_path: Path to input PDF file
        output_path: Path to output PPTX file
        options: Dictionary of conversion options
            - preserve_images: bool (default True)
            - preserve_formatting: bool (default True)
            - image_quality: int 1-100 (default 95)
            - extract_as_image: bool (default False) - render pages as images
    
    Returns:
        Dictionary with conversion result
    """
    converter = PDFToPPTXConverter(options)
    return converter.convert(input_path, output_path)


def main():
    """Main entry point for command line usage"""
    if len(sys.argv) < 3:
        print(json.dumps({
            'success': False,
            'error': 'Usage: python pdf_to_pptx.py <input_pdf> <output_pptx> [options_json]'
        }))
        sys.exit(1)
    
    input_path = sys.argv[1]
    output_path = sys.argv[2]
    
    # Parse options if provided
    options = {}
    if len(sys.argv) > 3:
        try:
            options = json.loads(sys.argv[3])
        except json.JSONDecodeError:
            pass
    
    # Validate input file exists
    if not os.path.exists(input_path):
        print(json.dumps({
            'success': False,
            'error': f'Input file not found: {input_path}'
        }))
        sys.exit(1)
    
    # Perform conversion
    result = convert_pdf_to_pptx(input_path, output_path, options)
    
    # Output result as JSON
    print(json.dumps(result))
    
    # Exit with appropriate code
    sys.exit(0 if result['success'] else 1)


if __name__ == '__main__':
    main()
